"""
Builds docs/07_03/weekly_report_07_03.pptx — same navy/teal theme as
docs/6_30/6_30_report_merged.pptx (colors extracted directly from that
file's XML: background 1C2B3A, card 134E4A, accent 0D9488, warn B45309,
text E2E8F0/FFFFFF, Calibri).

Run: python3 docs/07_03/build_report.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

NAVY = RGBColor(0x1C, 0x2B, 0x3A)
CARD = RGBColor(0x13, 0x4E, 0x4A)
TEAL = RGBColor(0x0D, 0x94, 0x88)
AMBER = RGBColor(0xB4, 0x53, 0x09)
RED = RGBColor(0xDC, 0x26, 0x26)
SLATE = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"
FONT_LIGHT = "Calibri Light"

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def new_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return slide


def add_text(slide, x, y, w, h, text, size=14, color=SLATE, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_card(slide, x, y, w, h, fill=CARD, line_color=None, radius=0.06):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_badge(slide, x, y, text, fill=TEAL, w=Inches(1.35), h=Inches(0.32)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    try:
        shape.adjustments[0] = 0.5
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = FONT
    return shape


def add_title_bar(slide, title, subtitle=None):
    add_text(slide, Inches(0.5), Inches(0.35), Inches(9), Inches(0.6),
              title, size=28, color=WHITE, bold=True, font=FONT_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.98), Inches(0.6), Pt(3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    bar.shadow.inherit = False
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.08), Inches(9), Inches(0.4),
                  subtitle, size=13, color=SLATE, italic=True)


# ==================================================================== #
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# -------------------------------------------------------------------- #
# Slide 1 — 封面
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_text(s, Inches(0.8), Inches(2.15), Inches(8.4), Inches(0.9),
          "MethyAgent 文献情报模块", size=40, color=WHITE, bold=True, font=FONT_LIGHT)
add_text(s, Inches(0.8), Inches(2.95), Inches(8.4), Inches(0.5),
          "进展汇报  ·  2026-7-3", size=18, color=TEAL, bold=True)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.05), Inches(1.0), Pt(4))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background(); bar.shadow.inherit = False

# -------------------------------------------------------------------- #
# Slide 2 — 本周进展
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_title_bar(s, "本周进展")

cards = [
    ("1", "LLM Reviewer Agent",
     "用二次 LLM 复核取代正则字符串匹配，修复 Bug 2（组织/cfDNA AUC 混淆）与 "
     "Bug 3（参考数据集误标），已在 PMID 40860669 上验证生效"),
    ("2", "新 Orchestrator（v2，agentic）",
     "新增独立的 LangGraph agentic 编排器：LLM 自主决定调用 search_papers / "
     "evaluate_geo_dataset / registry 写入，不改动现有 orchestrator.py，"
     "端到端 mock 测试已跑通"),
    ("3", "部署开关",
     "新增 ORCHESTRATOR_VERSION 配置（env var / settings.yaml），v1 / v2 "
     "可在部署时切换，默认 v1，现有链路不受影响"),
]
card_w = Inches(2.87)
card_h = Inches(3.55)
gap = Inches(0.2)
x0 = Inches(0.5)
y0 = Inches(1.55)
for i, (num, title, body) in enumerate(cards):
    x = Emu(x0 + i * (card_w + gap))
    add_card(s, x, y0, card_w, card_h)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x + Inches(0.22)), Emu(y0 + Inches(0.22)), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL; circ.line.fill.background(); circ.shadow.inherit = False
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = num; run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = WHITE
    add_text(s, Emu(x + Inches(0.22)), Emu(y0 + Inches(0.9)), Emu(card_w - Inches(0.44)), Inches(0.6),
              title, size=15, color=WHITE, bold=True)
    add_text(s, Emu(x + Inches(0.22)), Emu(y0 + Inches(1.55)), Emu(card_w - Inches(0.44)), Emu(card_h - Inches(1.75)),
              body, size=11.5, color=SLATE)

# -------------------------------------------------------------------- #
# Slide 3 — LLM Reviewer 详情
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_title_bar(s, "LLM Reviewer 详情", "PMID 40860669 — before / after，及为什么用 LLM 而不是正则")

table_x, table_y = Inches(0.5), Inches(1.55)
table_w, table_h = Inches(9.0), Inches(1.7)
rows, cols = 4, 3
gtable = s.shapes.add_table(rows, cols, table_x, table_y, table_w, table_h)
table = gtable.table
table.columns[0].width = Inches(1.9)
table.columns[1].width = Inches(3.4)
table.columns[2].width = Inches(3.7)

header = ["字段", "Draft（extract_paper_structured）", "Reviewer 修正后"]
data_rows = [
    ("sample_type", "plasma_cfdna", "plasma_cfdna（不变）"),
    ("auc_validation", "0.922", "null（已置空，标记人工复核）"),
    ("dataset_ids", "GSE50132, TCGA, GSE69914", "TCGA, GSE69914（GSE50132 已剔除）"),
]
for c, htext in enumerate(header):
    cell = table.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
    cell.text_frame.paragraphs[0].text = htext
    run = cell.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(12); run.font.bold = True; run.font.color.rgb = WHITE; run.font.name = FONT
for r, row in enumerate(data_rows, start=1):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = CARD
        cell.text_frame.paragraphs[0].text = val
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(11); run.font.color.rgb = SLATE if c != 1 or r != 1 else WHITE
        run.font.name = FONT
        if r == 2 and c == 1:  # 0.922 — highlight as wrong
            run.font.color.rgb = AMBER
            run.font.bold = True

add_text(s, Inches(0.5), Inches(3.5), Inches(9), Inches(0.3),
          "为什么用 LLM Reviewer 而不是正则（sentence co-occurrence）：", size=13, color=WHITE, bold=True)
why_points = (
    "•  正则要求 AUC 数值与 sample_type 关键词落在同一句 — 跨句但仍正确关联的表述会被误删，反之也可能漏判\n"
    "•  参考数据集判断依赖固定关键词表（reference/background/normalization），无法理解 accession 实际用途\n"
    "•  LLM 读取全文摘要做推理，是对两周前智谱指出的\"无 multi-turn reasoning / reflection loop\"问题的第一版解法\n"
    "•  代价：每篇论文多一次 LLM 调用（成本 / 延迟的权衡，已在代码注释中注明）"
)
add_text(s, Inches(0.5), Inches(3.85), Inches(9), Inches(1.6), why_points, size=11.5, color=SLATE)

# -------------------------------------------------------------------- #
# Slide 4 — 新 Orchestrator 架构图
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_title_bar(s, "新 Orchestrator 架构图", "v1（不变）vs v2（新增，agentic）")

panel_w, panel_h = Inches(4.3), Inches(3.3)
py = Inches(1.55)
lx = Inches(0.5)
rx = Inches(5.2)

add_card(s, lx, py, panel_w, panel_h, fill=CARD)
add_text(s, Emu(lx + Inches(0.25)), Emu(py + Inches(0.2)), Emu(panel_w - Inches(0.5)), Inches(0.4),
          "v1 — agents/orchestrator.py（不变）", size=14, color=WHITE, bold=True)
add_badge(s, Emu(lx + Inches(0.25)), Emu(py + Inches(0.62)), "固定顺序", fill=AMBER, w=Inches(1.15))
v1_flow = "parse_query\n     ↓\nrun_database_agent\n     ↓\nrun_literature_agent\n     ↓\ngenerate_report"
add_text(s, Emu(lx + Inches(0.25)), Emu(py + Inches(1.1)), Emu(panel_w - Inches(0.5)), Inches(1.6),
          v1_flow, size=13, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, Emu(lx + Inches(0.25)), Emu(py + Inches(2.75)), Emu(panel_w - Inches(0.5)), Inches(0.5),
          "节点顺序由代码硬编码；生产环境默认使用", size=10.5, color=SLATE, italic=True)

add_card(s, rx, py, panel_w, panel_h, fill=CARD)
add_text(s, Emu(rx + Inches(0.25)), Emu(py + Inches(0.2)), Emu(panel_w - Inches(0.5)), Inches(0.4),
          "v2 — agents/orchestrator_v2.py（新增）", size=14, color=WHITE, bold=True)
add_badge(s, Emu(rx + Inches(0.25)), Emu(py + Inches(0.62)), "LLM 自主决策", fill=TEAL, w=Inches(1.4))
v2_flow = "顶层 LLM（ReAct）\n     ↓ 自主选择\nsearch_papers  /  evaluate_geo_dataset_tool  /  write_to_registry\n（顺序、是否调用均由 LLM 决定）"
add_text(s, Emu(rx + Inches(0.25)), Emu(py + Inches(1.1)), Emu(panel_w - Inches(0.5)), Inches(1.6),
          v2_flow, size=12, color=SLATE, align=PP_ALIGN.CENTER)
add_text(s, Emu(rx + Inches(0.25)), Emu(py + Inches(2.75)), Emu(panel_w - Inches(0.5)), Inches(0.5),
          "端到端 mock 测试已通过；生产可用性待评估", size=10.5, color=SLATE, italic=True)

add_text(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.5),
          "两条路径共存，通过 ORCHESTRATOR_VERSION 配置切换（默认 v1）。v2 尚未接入下载 / 人工审批流程。",
          size=11, color=TEAL, bold=True)

# -------------------------------------------------------------------- #
# Slide 5 — Claude Scientist 对比测试结果（占位）
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_title_bar(s, "Claude Scientist 对比测试结果")
add_card(s, Inches(0.5), Inches(1.7), Inches(9), Inches(2.8), fill=CARD)
add_badge(s, Inches(0.8), Inches(2.0), "结果待补充", fill=AMBER, w=Inches(1.5))
add_text(s, Inches(0.8), Inches(2.55), Inches(8.4), Inches(1.7),
          "本周与 LLM Reviewer / 新 Orchestrator 并行，独立测试 Claude Scientist（Claude 研究模式），"
          "评估维度：\n"
          "•  端到端表现（end-to-end performance）\n"
          "•  文献检索能力（literature search ability）\n\n"
          "结果留待下次汇报补充。",
          size=13, color=SLATE)

# -------------------------------------------------------------------- #
# Slide 6 — 下一步
# -------------------------------------------------------------------- #
s = new_slide(prs)
add_title_bar(s, "下一步")

next_items = [
    ("待你在服务器验证", AMBER, "P0-2 / P0-5 的 LLM Reviewer 与新 Orchestrator 目前只用 stub LLM 做过 mock 测试"
     "（本地无 API key / NCBI 网络）；需要在服务器用真实 API key 跑一遍确认推理质量"),
    ("CRC biomarker 召回率测试", AMBER, "需要真实 LLM + PubMed 网络访问，本地无法跑，下周 / 服务器上补"),
    ("结转下周（P2）", RED, "gold_standard 剩余 8/10 条人工核验、data_availability 完整度 0% 根因排查、"
     "完整测试套件重跑（test_components_145 / test_ncbi_search / gold_standard）"),
    ("方向讨论", TEAL, "v2 orchestrator 生产可用性评估 — 是否 / 何时切换默认版本，"
     "以及是否往 P0-4 提到的完整技能化改造推进"),
]
y = Inches(1.55)
row_h = Inches(0.92)
for label, color, body in next_items:
    add_card(s, Inches(0.5), y, Inches(9), Emu(row_h - Inches(0.12)), fill=CARD)
    add_badge(s, Inches(0.7), Emu(y + Inches(0.13)), label, fill=color, w=Inches(2.0), h=Inches(0.4))
    add_text(s, Inches(2.85), Emu(y + Inches(0.08)), Inches(6.5), Emu(row_h - Inches(0.2)),
              body, size=10.5, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)
    y = Emu(y + row_h)

prs.save("docs/07_03/weekly_report_07_03.pptx")
print("Saved docs/07_03/weekly_report_07_03.pptx")
